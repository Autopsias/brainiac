"""PPTX handler — python-pptx. One `## Slide N` section per slide; text frames
+ tables (Markdown tables, headers retained)."""
from __future__ import annotations

from pathlib import Path

from .base import ExtractResult, Handler, density_gate, ocr_available, ocr_image
from .tables import rows_to_markdown

try:
    from pptx import Presentation
    _HAS_PPTX = True
except ImportError:  # pragma: no cover
    _HAS_PPTX = False

MAX_PPTX_BYTES = 150 * 1024 * 1024


def _ocr_pictures(slide: object) -> str:
    """OCR the pictures on one slide. A deck exported as one rendered image
    per slide carries ALL its text inside those pictures — the same problem a
    scanned PDF has, and the reference vault held a two-slide deck whose
    entire cost/timeline comparison was invisible for that reason. Returns ""
    when the slide has no picture or OCR reads nothing."""
    import io

    texts = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        blob = getattr(getattr(shape, "image", None), "blob", None)
        if blob is None:
            continue
        try:
            from PIL import Image

            with Image.open(io.BytesIO(blob)) as img:
                text, _ = ocr_image(img)
        except Exception:
            continue
        if text:
            texts.append(text)
    return "\n\n".join(texts).strip()


class PptxHandler(Handler):
    extensions = (".pptx",)
    dependency_name = "python-pptx"

    @classmethod
    def available(cls) -> bool:
        return _HAS_PPTX

    @classmethod
    def extract(cls, path: Path) -> ExtractResult:
        if not _HAS_PPTX:
            return ExtractResult.quarantine("missing_dependency:python-pptx")
        try:
            size = path.stat().st_size
        except OSError:
            size = 0
        if size > MAX_PPTX_BYTES:
            return ExtractResult.quarantine("file_too_large")
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            return ExtractResult.quarantine(
                "pptx_extraction_error", warnings=[f"{type(exc).__name__}: {exc}"]
            )

        sections: list[str] = []
        slide_count = 0
        ocr_slides: list[int] = []
        can_ocr = ocr_available()
        try:
            for i, slide in enumerate(prs.slides, start=1):
                slide_count = i
                lines = [f"## Slide {i}\n"]
                native = False
                for shape in slide.shapes:
                    if shape.has_table:
                        tbl = shape.table
                        rows = [[c.text for c in row.cells] for row in tbl.rows]
                        lines.append(rows_to_markdown(rows))
                        native = True
                    elif shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            lines.append(text + "\n")
                            native = True
                if not native and can_ocr:
                    # No native text on this slide — read its pictures.
                    picture_text = _ocr_pictures(slide)
                    if picture_text:
                        ocr_slides.append(i)
                        lines.append(picture_text + "\n")
                sections.append("\n".join(lines))
        except Exception as exc:
            return ExtractResult.quarantine(
                "pptx_extraction_error", warnings=[f"{type(exc).__name__}: {exc}"]
            )

        body = "\n".join(sections)
        reason = density_gate(body)
        if reason:
            return ExtractResult.quarantine(reason, warnings=[
                "no native text on any slide, and "
                + ("OCR of the slide pictures read nothing" if can_ocr
                   else "no local OCR engine is installed (pytesseract + tesseract)")])
        warnings = []
        if ocr_slides:
            warnings.append(
                f"ocr_slides: {len(ocr_slides)}/{slide_count} read by local OCR")
        return ExtractResult(markdown=body, warnings=warnings,
                             metadata={"slide_count": slide_count,
                                       "ocr_slides": ocr_slides})
