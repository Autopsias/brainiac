"""ING-01/ING-02 — ingestion pipeline core + document handlers.

Covers: happy path per handler, quarantine (unhandled ext / encrypted PDF /
empty-text), degraded-deps reporting, immutability of archived originals
(create-exclusive + collision detection), duplicate-content idempotency,
concurrent-ingest safety, and VM refusal with zero side effects.

Offline + deterministic: HashEmbedder + BruteForceBackend, env-injected audit
key (mirrors tests/test_capture_path.py conventions).
"""
from __future__ import annotations

import hashlib
import io
import threading
from pathlib import Path

import pytest

from brain.core import BrainCore, RoleError
from brain.embed import HashEmbedder
from brain.index import BrainIndex
from brain.vectors import BruteForceBackend

pytest.importorskip("docx")
pytest.importorskip("pptx")
pytest.importorskip("openpyxl")
pytest.importorskip("pypdf")

import docx
import openpyxl
from pptx import Presentation
from pypdf import PdfWriter


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _mini_vault(tmp_path: Path) -> Path:
    v = tmp_path / "vault"
    (v / "brain" / "resources").mkdir(parents=True)
    (v / "raw").mkdir(parents=True)
    (v / "inbox").mkdir(parents=True)
    return v


def _host_core(tmp_path: Path, vault: Path) -> BrainCore:
    idx = BrainIndex(db_path=tmp_path / "idx.sqlite", backend=BruteForceBackend(),
                      embedder=HashEmbedder())
    idx.rebuild(vault)
    return BrainCore(vault=vault, index=idx, audit_log=tmp_path / "audit.jsonl", role="host")


LOREM = ("This is a genuinely long paragraph of real prose used as fixture "
          "content so the extraction-quality density gate passes easily. ") * 3


def _make_docx(path: Path, *, with_table: bool = True) -> None:
    doc = docx.Document()
    doc.add_heading("Fixture document", level=1)
    doc.add_paragraph(LOREM)
    if with_table:
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Name"
        table.rows[0].cells[1].text = "Value"
        table.rows[1].cells[0].text = "alpha"
        table.rows[1].cells[1].text = "1"
    doc.save(path)


def _make_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Fixture slide"
    body = slide.placeholders[1]
    body.text_frame.text = LOREM
    prs.save(path)


def _make_xlsx(path: Path) -> None:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["item", "qty", "total"])
    ws.append(["widget", 2, "=B2*10"])
    ws.append(["gadget", 3, "=B3*10"])
    wb.save(path)


def _make_pdf_with_text(path: Path, text: str = "Hello ingestion world. " * 10) -> None:
    """A minimal single-page PDF with a real text-showing content stream —
    built by hand (no reportlab dependency) so pypdf's text extraction has
    something genuine to find."""
    writer = PdfWriter()
    page = writer.add_blank_page(width=200, height=200)
    stream = f"BT /F1 10 Tf 10 150 Td ({text}) Tj ET".encode("latin-1", "replace")
    from pypdf.generic import DictionaryObject, NameObject, ArrayObject, ContentStream
    cs = ContentStream(None, writer)
    cs.set_data(stream)
    page[NameObject("/Contents")] = writer._add_object(cs)
    resources = DictionaryObject()
    font = DictionaryObject({
        NameObject("/Type"): NameObject("/Font"),
        NameObject("/Subtype"): NameObject("/Type1"),
        NameObject("/BaseFont"): NameObject("/Helvetica"),
    })
    fonts = DictionaryObject({NameObject("/F1"): writer._add_object(font)})
    resources[NameObject("/Font")] = fonts
    page[NameObject("/Resources")] = resources
    with open(path, "wb") as f:
        writer.write(f)


def _make_blank_pdf(path: Path) -> None:
    """A page with no text content stream at all — the scanned/no-text-layer case."""
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with open(path, "wb") as f:
        writer.write(f)


def _make_encrypted_pdf(path: Path) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.encrypt("s3cret")
    with open(path, "wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# happy path per handler
# ---------------------------------------------------------------------------

class TestHappyPath:
    def test_docx_ingest(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_docx(vault / "inbox" / "report.docx")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert len(res["processed"]) == 1, res
        entry = res["processed"][0]
        note = (vault / entry["note"]).read_text(encoding="utf-8")
        assert "classification: Internal" in note
        assert "type: source" in note
        assert "| Name | Value |" in note or "table-unparsed" in note
        archived = vault / entry["archived"]
        assert archived.is_file(), "original must be archived immutably"
        assert not (vault / "inbox" / "report.docx").exists(), "claimed original must leave the inbox"

    def test_pptx_ingest(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_pptx(vault / "inbox" / "deck.pptx")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert len(res["processed"]) == 1, res
        note = (vault / res["processed"][0]["note"]).read_text(encoding="utf-8")
        assert "Fixture slide" in note or "Slide 1" in note

    def test_xlsx_ingest_preserves_table_and_formula(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_xlsx(vault / "inbox" / "sheet.xlsx")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert len(res["processed"]) == 1, res
        note = (vault / res["processed"][0]["note"]).read_text(encoding="utf-8")
        assert "| item | qty | total |" in note
        # openpyxl has no cached formula results for a workbook never opened
        # in Excel -> formula collapse degrades to the raw formula, tagged.
        assert "formula, uncomputed" in note

    def test_pdf_ingest_happy_path(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_pdf_with_text(vault / "inbox" / "doc.pdf")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert len(res["processed"]) == 1, res
        note = (vault / res["processed"][0]["note"]).read_text(encoding="utf-8")
        assert "Hello ingestion world" in note

    def test_text_ingest_happy_path(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "notes.txt").write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert len(res["processed"]) == 1, res


# ---------------------------------------------------------------------------
# quarantine paths
# ---------------------------------------------------------------------------

class TestQuarantine:
    def test_unhandled_extension_quarantines(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "mystery.xyz").write_bytes(b"unknown format")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert len(res["quarantined"]) == 1
        assert res["quarantined"][0]["reason"] == "no_handler_for_extension"
        assert (vault / "inbox" / "_quarantine" / "no_handler_for_extension" / "mystery.xyz").exists()
        assert not (vault / "raw" / "mystery.xyz").exists()

    def test_encrypted_pdf_quarantines_never_signs_garbage(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_encrypted_pdf(vault / "inbox" / "locked.pdf")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert res["processed"] == []
        assert len(res["quarantined"]) == 1
        assert res["quarantined"][0]["reason"] == "pdf_encrypted"
        assert not any((vault / "raw").glob("*.md"))

    def test_no_text_layer_pdf_quarantines(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_blank_pdf(vault / "inbox" / "scanned.pdf")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert res["processed"] == []
        assert res["quarantined"][0]["reason"] == "pdf_no_text_layer"

    def test_empty_text_docx_quarantines_quality_gate(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        doc = docx.Document()
        doc.add_paragraph("hi")  # far below the density gate
        doc.save(vault / "inbox" / "empty.docx")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert res["processed"] == []
        assert res["quarantined"][0]["reason"] == "empty_or_low_text_density"

    def test_malformed_office_file_quarantines_not_crashes(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "broken.docx").write_bytes(b"not a real zip/docx at all")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()  # must not raise
        assert res["quarantined"][0]["reason"] == "docx_extraction_error"

    def test_dry_run_makes_zero_filesystem_changes(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_docx(vault / "inbox" / "report.docx")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone(dry_run=True)
        assert res["processed"][0]["would_write"] is True
        assert (vault / "inbox" / "report.docx").exists(), "dry-run must not move the original"
        assert not any((vault / "raw").glob("*.md")), "dry-run must not write anything"


# ---------------------------------------------------------------------------
# degraded-deps path
# ---------------------------------------------------------------------------

class TestDegradedDeps:
    def test_missing_dependency_reports_and_quarantines(self, tmp_path, audit_key_env, monkeypatch):
        from brain.ingest.handlers import pdf as pdf_mod

        monkeypatch.setattr(pdf_mod, "_HAS_PYPDF", False)
        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "doc.pdf").write_bytes(b"%PDF-1.4 fake")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert res["quarantined"][0]["reason"] == "missing_dependency:pypdf"

    def test_capability_report_reflects_probe(self):
        from brain.ingest import capability_report

        report = capability_report()
        assert ".pdf" in report and ".docx" in report and ".xlsx" in report and ".pptx" in report
        assert all("available" in v for v in report.values())


# ---------------------------------------------------------------------------
# immutability + collisions (HARDENED:codex / codex-verify-r1)
# ---------------------------------------------------------------------------

class TestImmutabilityAndIdempotency:
    def test_duplicate_content_is_idempotent_noop(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_docx(vault / "inbox" / "first.docx")
        core = _host_core(tmp_path, vault)
        res1 = core.ingest_dropzone()
        assert len(res1["processed"]) == 1
        first_id = res1["processed"][0]["id"]

        # Re-drop IDENTICAL bytes under a different filename.
        import shutil
        shutil.copyfile(vault / "raw" / "originals" / f"{first_id}" / "first.docx",
                        vault / "inbox" / "second.docx")
        res2 = core.ingest_dropzone()
        assert res2["processed"] == [], "identical content must not be re-signed"
        assert len(res2["duplicates"]) == 1
        assert res2["duplicates"][0]["existing_id"] == first_id
        assert not (vault / "inbox" / "second.docx").exists()

    def test_archive_write_is_create_exclusive_same_sha_noop(self, tmp_path, audit_key_env):
        from brain.ingest.pipeline import _create_exclusive_or_collision

        target = tmp_path / "originals" / "x.bin"
        data = b"hello world"
        assert _create_exclusive_or_collision(target, data) == "written"
        assert _create_exclusive_or_collision(target, data) == "idempotent"

    def test_archive_write_collision_on_different_bytes_never_overwrites(self, tmp_path):
        from brain.ingest.pipeline import _create_exclusive_or_collision

        target = tmp_path / "originals" / "x.bin"
        assert _create_exclusive_or_collision(target, b"hello world") == "written"
        status = _create_exclusive_or_collision(target, b"DIFFERENT BYTES")
        assert status == "collision"
        assert target.read_bytes() == b"hello world", "collision must never clobber the existing archive"

    def test_archive_write_reuses_known_sha_on_collision_path(self, tmp_path):
        """SOAK-01: the caller-known sha256 must be honored on the
        collision-check path instead of re-hashing `data` a second time."""
        import hashlib
        from brain.ingest.pipeline import _create_exclusive_or_collision, _sha256_bytes

        target = tmp_path / "originals" / "x.bin"
        data = b"hello world"
        known_sha = _sha256_bytes(data)
        assert _create_exclusive_or_collision(target, data, known_sha=known_sha) == "written"

        # Re-ingest identical bytes with the correct known_sha: idempotent.
        assert _create_exclusive_or_collision(target, data, known_sha=known_sha) == "idempotent"

        # A deliberately WRONG known_sha must not be trusted blindly for
        # bytes that are actually identical to what's on disk — the fast
        # bytes-equality check must win over any caller-supplied digest.
        wrong_sha = hashlib.sha256(b"not the real content").hexdigest()
        assert _create_exclusive_or_collision(target, data, known_sha=wrong_sha) == "idempotent"

        # Different bytes + correct known_sha for the new content: collision,
        # and the known_sha must be used rather than recomputed (behavior
        # identical to the no-known_sha path).
        different = b"DIFFERENT BYTES"
        known_sha_different = _sha256_bytes(different)
        status = _create_exclusive_or_collision(target, different, known_sha=known_sha_different)
        assert status == "collision"
        assert target.read_bytes() == data


# ---------------------------------------------------------------------------
# concurrency (HARDENED:codex)
# ---------------------------------------------------------------------------

class TestConcurrency:
    def test_concurrent_ingest_never_double_processes(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_docx(vault / "inbox" / "race.docx")
        core = _host_core(tmp_path, vault)

        results = []

        def _run():
            results.append(core.ingest_dropzone())

        t1 = threading.Thread(target=_run)
        t2 = threading.Thread(target=_run)
        t1.start(); t2.start()
        t1.join(); t2.join()

        total_processed = sum(len(r["processed"]) for r in results)
        assert total_processed == 1, "exactly one thread must win the claim and process the file"
        notes = list((vault / "raw").glob("*.md"))
        assert len(notes) == 1

    def test_stale_processing_leftover_is_swept_and_reingested(self, tmp_path, audit_key_env):
        """A file stranded in inbox/_processing/ by a crashed prior process
        (claimed via os.rename, then the process died before it could
        unlink/promote or move it back) must not be lost forever — the NEXT
        run_ingest call sweeps it back to the inbox root and re-processes it."""
        import os as _os
        import time as _time

        from brain.ingest import pipeline as P

        vault = _mini_vault(tmp_path)
        core = _host_core(tmp_path, vault)

        processing = vault / "inbox" / P.PROCESSING_DIRNAME
        processing.mkdir(parents=True)
        stuck = processing / "crashed.docx"
        _make_docx(stuck)
        # Simulate the leftover being old enough to be a genuine stale claim
        # (not one currently being processed by a live, concurrent run).
        old = _time.time() - (P.STALE_PROCESSING_SECONDS + 60)
        _os.utime(stuck, (old, old))

        report = core.ingest_dropzone()

        assert len(report["processed"]) == 1
        assert report["processed"][0]["file"] == "crashed.docx"
        assert not stuck.exists()
        assert not any(processing.iterdir())
        assert len(list((vault / "raw").glob("*.md"))) == 1

    def test_recent_processing_leftover_is_not_swept(self, tmp_path, audit_key_env):
        """A FRESH _processing/ entry (within the staleness window) must be
        left alone — it may belong to another live process actively
        extracting it; sweeping it would break the atomic-claim guarantee."""
        from brain.ingest import pipeline as P

        vault = _mini_vault(tmp_path)
        core = _host_core(tmp_path, vault)

        processing = vault / "inbox" / P.PROCESSING_DIRNAME
        processing.mkdir(parents=True)
        live = processing / "live.docx"
        _make_docx(live)  # fresh mtime, well within the staleness window

        report = core.ingest_dropzone()

        assert report["processed"] == []
        assert live.exists(), "a recent claim must not be swept out from under a live process"

    def test_exception_mid_processing_moves_claim_back_for_next_run(self, tmp_path, audit_key_env, monkeypatch):
        """If write_note (or any per-file step) raises, the claimed file must
        be moved back to the inbox root — not stranded in _processing/ — so a
        subsequent run_ingest picks it up and retries.

        C2 (rework): this must NOT propagate out of ingest_dropzone — a single
        poison file re-raising on every call used to abort `brain sync`
        (and thus index reconciliation + snapshot publish) forever. The
        failure is recorded as a retryable "skipped" entry instead."""
        from brain.ingest import pipeline as P

        vault = _mini_vault(tmp_path)
        _make_docx(vault / "inbox" / "boom.docx")
        core = _host_core(tmp_path, vault)

        monkeypatch.setattr(
            core, "write_note",
            lambda *a, **kw: (_ for _ in ()).throw(RuntimeError("simulated crash")),
        )
        report = core.ingest_dropzone()  # must NOT raise
        assert report["processed"] == []
        assert any(s["file"] == "boom.docx" for s in report["skipped"]), report

        assert not (vault / "inbox" / "_processing").exists() or \
            not any((vault / "inbox" / P.PROCESSING_DIRNAME).iterdir())
        assert (vault / "inbox" / "boom.docx").exists(), \
            "claim must be moved back to inbox root on exception, not stranded"

        # Un-break write_note and confirm the next run succeeds normally.
        monkeypatch.undo()
        report = core.ingest_dropzone()
        assert len(report["processed"]) == 1

    def test_oversize_file_quarantined_without_full_read(self, tmp_path, audit_key_env, monkeypatch):
        """The size cap must be checked via stat() BEFORE read_bytes() so a
        pathological huge file is never loaded fully into memory."""
        from brain.ingest import pipeline as P

        vault = _mini_vault(tmp_path)
        oversize = vault / "inbox" / "huge.docx"
        oversize.write_bytes(b"not-a-real-docx")  # content is irrelevant; size gate short-circuits first
        core = _host_core(tmp_path, vault)

        monkeypatch.setattr(P, "MAX_INGEST_BYTES", 1)  # any non-empty file exceeds this
        read_calls = {"n": 0}
        real_read_bytes = Path.read_bytes

        def _tracking_read_bytes(self):
            if self.name == "huge.docx":
                read_calls["n"] += 1
            return real_read_bytes(self)

        monkeypatch.setattr(Path, "read_bytes", _tracking_read_bytes)

        report = core.ingest_dropzone()

        assert read_calls["n"] == 0, "oversize file must be rejected via stat() before any read_bytes()"
        assert len(report["quarantined"]) == 1
        assert report["quarantined"][0]["reason"] == "file_too_large"

    def test_processed_entry_carries_real_classification(self, tmp_path, audit_key_env):
        """cli.py's ingest egress gate must see the promoted note's ACTUAL
        classification, not a fabricated constant."""
        vault = _mini_vault(tmp_path)
        _make_docx(vault / "inbox" / "doc.docx")
        core = _host_core(tmp_path, vault)

        report = core.ingest_dropzone()

        assert len(report["processed"]) == 1
        entry = report["processed"][0]
        assert "classification" in entry
        note_path = vault / entry["note"]
        from brain import frontmatter as fm
        meta, _ = fm.parse_text(note_path.read_text(encoding="utf-8"))
        assert entry["classification"] == meta["classification"]


# ---------------------------------------------------------------------------
# VM refusal (zero side effects before refusal)
# ---------------------------------------------------------------------------

class TestVmRefusal:
    def test_vm_ingest_refused_with_zero_side_effects(self, tmp_path):
        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "report.docx").write_bytes(b"placeholder")
        snap_idx = BrainIndex(db_path=tmp_path / "snap.sqlite", backend=BruteForceBackend(),
                               embedder=HashEmbedder(), read_only=True)
        core = BrainCore(vault=vault, index=snap_idx, role="vm")

        with pytest.raises(RoleError):
            core.ingest_dropzone()

        # Zero side effects: no processing/quarantine dirs, original untouched.
        assert not (vault / "inbox" / "_processing").exists()
        assert not (vault / "inbox" / "_quarantine").exists()
        assert (vault / "inbox" / "report.docx").exists()
        assert not any((vault / "raw").glob("*.md"))

    def test_vm_ingest_never_resolves_signing_key(self, tmp_path, monkeypatch):
        import brain.audit as audit_mod

        called = {"n": 0}
        monkeypatch.setattr(
            audit_mod, "resolve_signing_key",
            lambda *a, **kw: (_ for _ in ()).throw(AssertionError("must not be called")),
        )
        vault = _mini_vault(tmp_path)
        snap_idx = BrainIndex(db_path=tmp_path / "snap.sqlite", backend=BruteForceBackend(),
                               embedder=HashEmbedder(), read_only=True)
        core = BrainCore(vault=vault, index=snap_idx, role="vm")
        with pytest.raises(RoleError):
            core.ingest_dropzone()


# ---------------------------------------------------------------------------
# sync() fold — ADR-0003 Ruling 1 amendment (drain-on-invoke, not nightly-only)
# ---------------------------------------------------------------------------

class TestSyncFold:
    def test_sync_drains_the_inbox(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        _make_docx(vault / "inbox" / "report.docx")
        core = _host_core(tmp_path, vault)
        res = core.sync(drain=True, publish=False)
        assert "ingest" in res
        assert len(res["ingest"]["processed"]) == 1
        note_id = res["ingest"]["processed"][0]["id"]
        assert core.get(note_id) is not None, "sync must reconcile the index after ingest"

    def test_sync_ingest_is_cheap_and_idempotent_when_inbox_empty(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        core = _host_core(tmp_path, vault)
        res = core.sync(drain=True)
        assert res["ingest"]["processed"] == []

    def test_sync_never_aborts_index_reconcile_when_ingest_blows_up(self, tmp_path, audit_key_env, monkeypatch):
        """C2 (core.sync backstop): ingest_dropzone() ran BEFORE index.sync()
        with no try/except — any exception escaping run_ingest's own per-file
        retry machinery (e.g. a manifest/failures-file I/O error) aborted
        index reconciliation and snapshot publication on every subsequent
        sync. core.sync() must catch it and still reconcile the index."""
        vault = _mini_vault(tmp_path)
        core = _host_core(tmp_path, vault)

        monkeypatch.setattr(
            core, "ingest_dropzone",
            lambda *a, **kw: (_ for _ in ()).throw(RuntimeError("simulated ingest blowup")),
        )
        res = core.sync(drain=True)  # must NOT raise
        assert "error" in res["ingest"]
        assert res["ingest"]["processed"] == []
        # the index reconcile must still have run (not short-circuited).
        assert "mode" in res and "chunks" in res


# ---------------------------------------------------------------------------
# S05 rework — C1: os.rename preserves mtime, so _claim must touch it
# ---------------------------------------------------------------------------

class TestClaimMtime:
    def test_claim_touches_mtime_so_stale_sweep_does_not_reclaim_a_live_claim(self, tmp_path):
        """A downloaded/copied source file commonly keeps an OLD mtime. Before
        the fix, _sweep_stale_processing (keyed on st_mtime) would see that
        old mtime immediately after the claim's os.rename (which PRESERVES
        mtime) and sweep a live, just-claimed file back to the inbox —
        exactly the double-process/duplicate-write race the claim mechanism
        exists to prevent."""
        import os as _os
        import time as _time

        from brain.ingest import pipeline as P

        vault = _mini_vault(tmp_path)
        src = vault / "inbox" / "old.txt"
        src.write_text("hello world", encoding="utf-8")
        old = _time.time() - (P.STALE_PROCESSING_SECONDS + 300)
        _os.utime(src, (old, old))

        processing_dir = vault / "inbox" / P.PROCESSING_DIRNAME
        claimed = P._claim(src, processing_dir)
        assert claimed is not None

        # Simulate a concurrent drain running the stale-sweep immediately
        # after the claim.
        P._sweep_stale_processing(
            processing_dir, vault / "inbox", vault=vault,
            quarantine_dir=vault / "inbox" / P.QUARANTINE_DIRNAME, failures={},
        )

        assert claimed.exists(), (
            "a freshly-claimed file must not be swept as 'stale' merely "
            "because its ORIGINAL source mtime was old"
        )


# ---------------------------------------------------------------------------
# S05 rework — C2: a poison file must never abort the whole drain, but must
# eventually be quarantined rather than retried forever.
# ---------------------------------------------------------------------------

class TestFailureIsolation:
    def test_poison_file_does_not_abort_the_run_or_later_files(self, tmp_path, audit_key_env, monkeypatch):
        from brain.ingest.handlers import text as text_mod

        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "aaa_poison.txt").write_text(LOREM, encoding="utf-8")
        (vault / "inbox" / "zzz_good.txt").write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        real_extract = text_mod.TextHandler.extract.__func__

        def _boom(cls, path):
            if path.name == "aaa_poison.txt":
                raise RuntimeError("simulated poison file")
            return real_extract(cls, path)

        monkeypatch.setattr(text_mod.TextHandler, "extract", classmethod(_boom))

        report = core.ingest_dropzone()  # must NOT raise

        assert len(report["processed"]) == 1, report
        assert report["processed"][0]["file"] == "zzz_good.txt", (
            "a poison file earlier in sort order must not prevent a "
            "later-alphabetical file from being processed in the SAME run"
        )
        assert any(s["file"] == "aaa_poison.txt" for s in report["skipped"]), report
        assert (vault / "inbox" / "aaa_poison.txt").exists(), \
            "poison file must be retried (moved back to inbox), not lost"

    def test_repeated_failure_quarantines_instead_of_retrying_forever(self, tmp_path, audit_key_env, monkeypatch):
        from brain.ingest import pipeline as P
        from brain.ingest.handlers import text as text_mod

        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "poison.txt").write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        def _always_boom(cls, path):
            raise RuntimeError("simulated poison file")

        monkeypatch.setattr(text_mod.TextHandler, "extract", classmethod(_always_boom))

        report = {}
        for _ in range(P.MAX_INGEST_FAILURES):
            report = core.ingest_dropzone()  # must never raise, on any attempt

        assert report.get("quarantined"), report
        assert report["quarantined"][0]["reason"] == "repeated_ingest_failure"
        bucket = vault / "inbox" / P.QUARANTINE_DIRNAME / "repeated_ingest_failure"
        assert list(bucket.glob("poison*")), \
            "poison file must land in _quarantine after MAX_INGEST_FAILURES attempts"
        assert not (vault / "inbox" / "poison.txt").exists()
        # one more drain must be a no-op (nothing left to retry/fail on).
        report2 = core.ingest_dropzone()
        assert report2["quarantined"] == []
        assert report2["skipped"] == []


# ---------------------------------------------------------------------------
# S05 rework — C3: _quarantine must never silently clobber a same-named entry
# ---------------------------------------------------------------------------

class TestQuarantineUniqueness:
    def test_two_same_named_quarantined_originals_both_survive(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        core = _host_core(tmp_path, vault)

        (vault / "inbox" / "mystery.xyz").write_bytes(b"first payload")
        report1 = core.ingest_dropzone()
        assert report1["quarantined"][0]["reason"] == "no_handler_for_extension"

        (vault / "inbox" / "mystery.xyz").write_bytes(b"second DIFFERENT payload")
        report2 = core.ingest_dropzone()
        assert report2["quarantined"][0]["reason"] == "no_handler_for_extension"

        bucket = vault / "inbox" / "_quarantine" / "no_handler_for_extension"
        payloads = {p.read_bytes() for p in bucket.glob("mystery*.xyz")}
        assert payloads == {b"first payload", b"second DIFFERENT payload"}, (
            "a second same-named quarantined file must never clobber the "
            "first — both are the only copy of their original content"
        )
        # each survivor must carry its OWN reason report, not a shared/lost one.
        reason_files = list(bucket.glob("mystery*.reason.txt"))
        assert len(reason_files) == 2


# ---------------------------------------------------------------------------
# S05 rework — C4: the inbox-exclusion check must anchor to the vault-relative
# TOP-LEVEL path segment, not an unanchored "/inbox/" substring match.
# ---------------------------------------------------------------------------

class TestInboxAnchoring:
    def test_nested_inbox_named_folder_note_is_still_indexed(self, tmp_path):
        from brain.notes import scan_vault

        vault = _mini_vault(tmp_path)
        nested = vault / "brain" / "resources" / "inbox"
        nested.mkdir(parents=True)
        (nested / "reading-list.md").write_text(
            "---\nid: reading-list\ntitle: \"Reading List\"\ntype: note\n"
            "classification: Internal\ncreated: 2026-07-05\nupdated: 2026-07-05\n"
            "---\n\nSome content under a folder merely NAMED inbox.\n",
            encoding="utf-8",
        )
        ids = {n.id for n in scan_vault(vault)}
        assert "reading-list" in ids, (
            "a note under a directory merely named 'inbox' (not the top-level "
            "drop zone) must not be silently dropped from the index"
        )

    def test_top_level_inbox_dropzone_is_still_excluded(self, tmp_path):
        from brain.notes import scan_vault

        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "draft.md").write_text(
            "---\nid: draft\ntitle: \"Draft\"\ntype: note\nclassification: Internal\n"
            "created: 2026-07-05\nupdated: 2026-07-05\n---\n\nDrop-zone content.\n",
            encoding="utf-8",
        )
        ids = {n.id for n in scan_vault(vault)}
        assert "draft" not in ids, "the real top-level inbox/ drop zone must still be excluded"

    def test_nested_inbox_named_folder_note_passes_validate(self, tmp_path, monkeypatch):
        import sys as _sys

        _sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "tools"))
        import validate as V

        vault = _mini_vault(tmp_path)
        (vault / "brain" / "index.md").write_text(
            "---\nid: index\ntitle: \"Index\"\ntype: index\nclassification: Internal\n"
            "created: 2026-07-05\nupdated: 2026-07-05\n---\n\nMap.\n", encoding="utf-8",
        )
        nested = vault / "brain" / "resources" / "inbox"
        nested.mkdir(parents=True)
        (nested / "reading-list.md").write_text(
            "---\nid: reading-list\ntitle: \"Reading List\"\ntype: note\n"
            "classification: Internal\ncreated: 2026-07-05\nupdated: 2026-07-05\n"
            "---\n\nContent under a folder merely named inbox.\n", encoding="utf-8",
        )
        monkeypatch.setattr(_sys, "argv", ["validate.py", str(vault)])
        rc = V.main()
        assert rc == 0, (V.errors, V.warnings)
        assert any(p.name == "reading-list.md" for p in V.iter_md(vault / "brain", vault))


# ---------------------------------------------------------------------------
# S05 rework — C5: raw/originals/ (archived ingestion originals) must never
# be scanned as a note or validated as one.
# ---------------------------------------------------------------------------

class TestArchivedOriginalsExcluded:
    def test_archived_frontmatterless_md_excluded_from_scan_and_validate(
        self, tmp_path, audit_key_env, monkeypatch,
    ):
        from brain.notes import scan_vault

        vault = _mini_vault(tmp_path)
        (vault / "brain" / "index.md").write_text(
            "---\nid: index\ntitle: \"Index\"\ntype: index\nclassification: Internal\n"
            "created: 2026-07-05\nupdated: 2026-07-05\n---\n\nMap.\n", encoding="utf-8",
        )
        (vault / "inbox" / "notes-plain.md").write_text(
            "# Just some plain markdown\n\n" + LOREM, encoding="utf-8",
        )
        core = _host_core(tmp_path, vault)
        report = core.ingest_dropzone()
        assert len(report["processed"]) == 1, report
        archived = vault / report["processed"][0]["archived"]
        assert archived.is_file()
        assert archived.parent.parent.name == "originals"

        scanned_paths = {n.path for n in scan_vault(vault)}
        assert archived not in scanned_paths, (
            "the archived ORIGINAL (frontmatter-less) must never be surfaced "
            "as a note by scan_vault"
        )

        import sys as _sys

        _sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "tools"))
        import validate as V

        monkeypatch.setattr(_sys, "argv", ["validate.py", str(vault)])
        rc = V.main()
        assert rc == 0, (
            "a frontmatter-less archived original must not break the "
            f"conventions gate: errors={V.errors}"
        )

    def test_archived_md_with_its_own_frontmatter_not_double_indexed(self, tmp_path, audit_key_env):
        from brain.notes import scan_vault

        vault = _mini_vault(tmp_path)
        content = (
            "---\nid: foreign-note\ntitle: \"Foreign\"\ntype: note\n"
            "classification: Internal\ncreated: 2026-07-01\nupdated: 2026-07-01\n"
            "---\n\n" + LOREM
        )
        (vault / "inbox" / "foreign.md").write_text(content, encoding="utf-8")
        core = _host_core(tmp_path, vault)
        report = core.ingest_dropzone()
        assert len(report["processed"]) == 1, report

        ids = [n.id for n in scan_vault(vault)]
        assert ids.count("foreign-note") == 0, (
            "the archived original (which happens to carry its OWN "
            "frontmatter id) must never surface as a duplicate/foreign-id "
            "raw note alongside the properly-promoted raw/<slug>.md"
        )


# ---------------------------------------------------------------------------
# S05 rework — C6: the dry-run preview path needs the same stat()-based size
# gate as the real ingest path, BEFORE any read_bytes()/extract().
# ---------------------------------------------------------------------------

class TestDryRunSizeGate:
    def test_dry_run_oversize_file_rejected_without_full_read(self, tmp_path, audit_key_env, monkeypatch):
        from brain.ingest import pipeline as P

        vault = _mini_vault(tmp_path)
        oversize = vault / "inbox" / "huge.txt"
        oversize.write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        monkeypatch.setattr(P, "MAX_INGEST_BYTES", 1)  # any non-empty file exceeds this
        read_calls = {"n": 0}
        real_read_bytes = Path.read_bytes

        def _tracking_read_bytes(self):
            if self.name == "huge.txt":
                read_calls["n"] += 1
            return real_read_bytes(self)

        monkeypatch.setattr(Path, "read_bytes", _tracking_read_bytes)

        report = core.ingest_dropzone(dry_run=True)

        assert read_calls["n"] == 0, \
            "dry-run must reject an oversize file via stat() before any read_bytes()"
        assert any(
            q["file"] == "huge.txt" and q["reason"] == "file_too_large"
            for q in report["quarantined"]
        ), report
        assert oversize.exists(), "dry-run must never move/touch the original file"


# ---------------------------------------------------------------------------
# S05 rework — C7: frontmatter must properly escape embedded quotes, and the
# archived-original filename component must be sanitized before it flows
# into the signed, immutable `origin:` value.
# ---------------------------------------------------------------------------

class TestFrontmatterEscaping:
    def test_build_frontmatter_escapes_embedded_quotes(self):
        from brain.ingest.pipeline import _build_frontmatter

        yaml = pytest.importorskip("yaml")
        meta = {
            "id": "x",
            "origin": 'weird: value with "quotes" and : a colon',
            "immutable": True,
        }
        text = _build_frontmatter(meta, "body text")
        block = text.split("---", 2)[1]
        parsed = yaml.safe_load(block)
        assert parsed["origin"] == 'weird: value with "quotes" and : a colon', (
            "an embedded double-quote must round-trip through real YAML, "
            "not silently corrupt the frontmatter block"
        )

    def test_hostile_filename_is_sanitized_and_note_stays_parseable(self, tmp_path, audit_key_env):
        from brain import frontmatter as fm

        vault = _mini_vault(tmp_path)
        hostile_name = 'report:"final".txt'
        (vault / "inbox" / hostile_name).write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        report = core.ingest_dropzone()
        assert len(report["processed"]) == 1, report

        archived_rel = report["processed"][0]["archived"]
        assert '"' not in archived_rel and ':' not in Path(archived_rel).name, (
            "the archived original's filename component must be sanitized "
            "of quote/colon chars before it flows into origin: / the archive path"
        )
        assert (vault / archived_rel).is_file()

        note_path = vault / report["processed"][0]["note"]
        text = note_path.read_text(encoding="utf-8")
        meta, _ = fm.parse_text(text)
        assert meta["id"] == report["processed"][0]["id"]

        yaml = pytest.importorskip("yaml")
        block = text.split("---", 2)[1]
        parsed = yaml.safe_load(block)
        assert isinstance(parsed, dict) and parsed["id"] == meta["id"]


# ---------------------------------------------------------------------------
# S05 rework — C8: sync --json must route the ingest sub-report's promoted
# entries through the SAME egress gate `ingest --json` already applies.
# ---------------------------------------------------------------------------

class TestSyncEgressGate:
    def test_sync_json_applies_egress_gate_to_ingest_processed(self, tmp_path, monkeypatch):
        import json as _json
        from contextlib import redirect_stdout

        from brain import cli
        from brain.core import BrainCore

        vault = tmp_path / "vault"
        (vault / "brain").mkdir(parents=True)
        (vault / "brain" / "index.md").write_text(
            "---\nid: index\ntitle: \"Index\"\ntype: index\nclassification: Internal\n"
            "created: 2026-07-05\nupdated: 2026-07-05\n---\n\nMap.\n", encoding="utf-8",
        )
        (vault / "raw").mkdir(parents=True)

        fake_result = {
            "mode": "incremental", "added": 0, "updated": 0, "deleted": 0, "unchanged": 0,
            "chunks": 0,
            "drain": {"promoted": 0, "skipped": 0},
            "ingest": {
                "processed": [
                    {"file": "secret.pdf", "id": "leak-id", "note": "raw/leak-id.md",
                     "classification": "Restricted", "warnings": []},
                ],
                "quarantined": [], "duplicates": [], "skipped": [], "dry_run": False,
            },
        }
        monkeypatch.setattr(BrainCore, "sync", lambda self, **kw: fake_result)
        monkeypatch.setenv("BRAIN_VAULT", str(vault))
        monkeypatch.setenv("BRAIN_INDEX_DIR", str(tmp_path / "idx"))

        buf = io.StringIO()
        with redirect_stdout(buf):
            code = cli.main(["sync", "--json"])
        assert code == 0
        out = _json.loads(buf.getvalue())

        assert out["ingest"]["processed"] == [], out
        assert out["ingest"]["egress"]["withheld"] == 1


# ---------------------------------------------------------------------------
# S05 rework — C9/C10: xlsx/docx quadratic-access rewrites must preserve
# behavior exactly (perf test optional per the finding; these are the
# correctness safety nets for the rewritten lookup logic).
# ---------------------------------------------------------------------------

class TestXlsxLockstepFormulaAlignment:
    def test_multiple_formula_columns_align_correctly_after_lockstep_rewrite(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["item", "qty", "total", "note"])
        ws.append(["widget", 2, "=B2*10", "plain"])
        ws.append(["gadget", 3, "=B3*10", '=CONCATENATE(A3,"-ok")'])
        wb.save(vault / "inbox" / "multi.xlsx")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert len(res["processed"]) == 1, res
        note = (vault / res["processed"][0]["note"]).read_text(encoding="utf-8")
        # both formula cells (different rows AND columns) must resolve to
        # their OWN formula text, not a neighbor's — this is exactly what a
        # row/column misalignment in the lockstep zip() rewrite would break.
        assert "=B2*10 (formula, uncomputed)" in note
        assert "=B3*10 (formula, uncomputed)" in note
        assert "CONCATENATE" in note
        assert "plain" in note


class TestDocxParagraphLookupCorrectness:
    def test_paragraph_table_document_order_preserved_after_dict_lookup_rewrite(self, tmp_path, audit_key_env):
        vault = _mini_vault(tmp_path)
        doc = docx.Document()
        doc.add_heading("Title One", level=1)
        doc.add_paragraph(LOREM)
        table = doc.add_table(rows=1, cols=2)
        table.rows[0].cells[0].text = "K"
        table.rows[0].cells[1].text = "V"
        doc.add_paragraph("Trailing paragraph after the table. " + LOREM)
        doc.save(vault / "inbox" / "multi.docx")
        core = _host_core(tmp_path, vault)
        res = core.ingest_dropzone()
        assert len(res["processed"]) == 1, res
        note = (vault / res["processed"][0]["note"]).read_text(encoding="utf-8")
        title_idx = note.index("Title One")
        table_idx = note.index("| K | V |") if "| K | V |" in note else note.index("table-unparsed")
        trailing_idx = note.index("Trailing paragraph")
        assert title_idx < table_idx < trailing_idx, (
            "document order (paragraph -> table -> paragraph) must be "
            "preserved by the element->paragraph dict lookup rewrite"
        )


# ---------------------------------------------------------------------------
# S05 extra fix pass — B1: a SYSTEMIC failure (no signing key, disk full) must
# never be counted against the per-file poison counter or quarantine a
# perfectly good file. Only a genuine per-file extraction/content defect may.
# ---------------------------------------------------------------------------

class TestSystemicVsPerFileFailure:
    def test_key_unavailable_does_not_poison_or_quarantine_any_file(
        self, tmp_path, audit_key_env, monkeypatch
    ):
        from brain.audit import KeyUnavailable
        from brain.ingest import pipeline as P

        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "a_first.txt").write_text(LOREM, encoding="utf-8")
        (vault / "inbox" / "b_second.txt").write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        def _boom(self, *a, **kw):
            raise KeyUnavailable("simulated signing-key outage")

        monkeypatch.setattr(BrainCore, "write_note", _boom)

        for _ in range(P.MAX_INGEST_FAILURES):
            report = core.ingest_dropzone()  # must NOT raise
            assert report["quarantined"] == [], (
                "a systemic outage must never quarantine a file", report
            )

        remaining = {p.name for p in (vault / "inbox").iterdir() if p.is_file()}
        assert remaining == {"a_first.txt", "b_second.txt"}, (
            "every file must still be sitting in the inbox for retry once "
            "the outage clears", remaining
        )
        quarantine_dir = vault / "inbox" / P.QUARANTINE_DIRNAME
        assert not quarantine_dir.exists() or not list(quarantine_dir.rglob("*.reason.txt"))
        assert P._load_failures(vault) == {}, (
            "a systemic (KeyUnavailable) failure must never bump the "
            "per-file poison counter"
        )

    def test_genuinely_poison_file_still_quarantines(self, tmp_path, audit_key_env, monkeypatch):
        """Companion to the systemic case above: a REAL per-file defect must
        still quarantine after MAX_INGEST_FAILURES — this is the existing
        behavior in TestFailureIsolation, re-asserted here to pin the
        systemic/per-file split doesn't accidentally make genuine poison
        un-quarantinable."""
        from brain.ingest import pipeline as P
        from brain.ingest.handlers import text as text_mod

        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "poison.txt").write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        def _always_boom(cls, path):
            raise RuntimeError("simulated poison file (per-file, not systemic)")

        monkeypatch.setattr(text_mod.TextHandler, "extract", classmethod(_always_boom))

        report = {}
        for _ in range(P.MAX_INGEST_FAILURES):
            report = core.ingest_dropzone()

        assert report["quarantined"], report
        assert report["quarantined"][0]["reason"] == "repeated_ingest_failure"
        bucket = vault / "inbox" / P.QUARANTINE_DIRNAME / "repeated_ingest_failure"
        assert list(bucket.glob("poison*"))


# ---------------------------------------------------------------------------
# S05 extra fix pass — B2: a filename carrying a control char (e.g. a literal
# newline) must not corrupt the signed frontmatter it flows into via origin:.
# ---------------------------------------------------------------------------

class TestControlCharFilenameSanitization:
    def test_newline_and_quote_in_filename_yields_parseable_frontmatter(
        self, tmp_path, audit_key_env
    ):
        from brain import frontmatter as fm

        vault = _mini_vault(tmp_path)
        hostile_name = 'evil\n"report".txt'
        (vault / "inbox" / hostile_name).write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        report = core.ingest_dropzone()
        assert len(report["processed"]) == 1, report

        archived_rel = report["processed"][0]["archived"]
        archived_name = Path(archived_rel).name
        assert "\n" not in archived_name and '"' not in archived_name, (
            "a control char / quote in the original filename must be "
            "stripped before it flows into the archive path / origin: value"
        )

        note_path = vault / report["processed"][0]["note"]
        text = note_path.read_text(encoding="utf-8")
        meta, _ = fm.parse_text(text)
        assert meta["id"] == report["processed"][0]["id"]
        assert "\n" not in meta.get("origin", "")

        yaml = pytest.importorskip("yaml")
        block = text.split("---", 2)[1]
        parsed = yaml.safe_load(block)
        assert isinstance(parsed, dict) and parsed["id"] == meta["id"], (
            "the frontmatter block must round-trip through real YAML even "
            "when the original filename carried a newline + a quote"
        )


# ---------------------------------------------------------------------------
# S05 extra fix pass — E3: a process-death (OOM/segfault) mid-extraction must
# not be re-swept forever — after enough stale-sweeps it must quarantine like
# any other repeated failure.
# ---------------------------------------------------------------------------

class TestProcessDeathPoisonLoop:
    def test_repeatedly_stale_swept_file_quarantines_instead_of_looping_forever(
        self, tmp_path
    ):
        import os as _os
        import time as _time

        from brain.ingest import pipeline as P

        vault = _mini_vault(tmp_path)
        processing_dir = vault / "inbox" / P.PROCESSING_DIRNAME
        quarantine_dir = vault / "inbox" / P.QUARANTINE_DIRNAME
        processing_dir.mkdir(parents=True)
        old = _time.time() - (P.STALE_PROCESSING_SECONDS + 300)

        stuck = processing_dir / "crashes-every-time.txt"
        stuck.write_text("payload that reliably kills the process", encoding="utf-8")
        _os.utime(stuck, (old, old))

        for attempt in range(1, P.MAX_INGEST_FAILURES + 1):
            failures = P._load_failures(vault)
            P._sweep_stale_processing(
                processing_dir, vault / "inbox", vault=vault,
                quarantine_dir=quarantine_dir, failures=failures,
            )
            if attempt < P.MAX_INGEST_FAILURES:
                # Swept back to the inbox root; simulate ANOTHER process
                # re-claiming it and dying mid-extraction before ITS next sweep.
                swept = vault / "inbox" / "crashes-every-time.txt"
                assert swept.exists(), f"attempt {attempt}: must be swept back, not lost"
                _os.rename(swept, stuck)
                _os.utime(stuck, (old, old))

        bucket = quarantine_dir / "repeated_ingest_failure"
        assert list(bucket.glob("crashes-every-time*")), (
            "a file that dies mid-extraction every time it's claimed must "
            "eventually quarantine rather than loop through the stale sweep "
            "forever"
        )
        assert P._load_failures(vault) == {}, (
            "the counter entry must be cleared once the file quarantines"
        )


# ---------------------------------------------------------------------------
# S05 extra fix pass — E4: the `duplicates` sub-list carries a real note id
# (`existing_id`) and must be routed through the SAME egress gate as
# `processed`, in both `ingest --json` and `sync --json`.
# ---------------------------------------------------------------------------

class TestDuplicateEgressGate:
    def test_duplicate_report_carries_existing_notes_classification(
        self, tmp_path, audit_key_env
    ):
        """Pipeline-level: the duplicate entry must thread through the
        EXISTING note's real classification (not the ingest-default
        'Internal'), so the CLI's egress gate has something real to filter on."""
        import shutil

        vault = _mini_vault(tmp_path)
        _make_docx(vault / "inbox" / "first.docx")
        core = _host_core(tmp_path, vault)
        res1 = core.ingest_dropzone()
        first_id = res1["processed"][0]["id"]

        # Simulate the existing note having since been reclassified (manifest
        # dedup keys on the ORIGINAL FILE bytes, not the note's current
        # frontmatter, so this must not affect duplicate DETECTION — only
        # what classification the duplicate-report entry carries).
        note_path = vault / "raw" / f"{first_id}.md"
        text = note_path.read_text(encoding="utf-8")
        note_path.write_text(
            text.replace("classification: Internal", "classification: Restricted"),
            encoding="utf-8",
        )

        shutil.copyfile(
            vault / "raw" / "originals" / f"{first_id}" / "first.docx",
            vault / "inbox" / "second.docx",
        )
        res2 = core.ingest_dropzone()
        assert len(res2["duplicates"]) == 1
        assert res2["duplicates"][0]["classification"] == "Restricted"

    def test_ingest_json_filters_above_tier_duplicate(self, tmp_path, monkeypatch):
        import json as _json
        from contextlib import redirect_stdout

        from brain import cli

        vault = tmp_path / "vault"
        (vault / "brain").mkdir(parents=True)
        (vault / "brain" / "index.md").write_text(
            "---\nid: index\ntitle: \"Index\"\ntype: index\nclassification: Internal\n"
            "created: 2026-07-05\nupdated: 2026-07-05\n---\n\nMap.\n", encoding="utf-8",
        )
        (vault / "raw").mkdir(parents=True)

        fake_result = {
            "processed": [], "quarantined": [], "skipped": [], "dry_run": False,
            "duplicates": [
                {"file": "dup.pdf", "existing_id": "leak-id", "classification": "Restricted"},
            ],
        }
        monkeypatch.setattr(BrainCore, "ingest_dropzone", lambda self, **kw: fake_result)
        monkeypatch.setenv("BRAIN_VAULT", str(vault))
        monkeypatch.setenv("BRAIN_INDEX_DIR", str(tmp_path / "idx"))

        buf = io.StringIO()
        with redirect_stdout(buf):
            code = cli.main(["ingest", "--json"])
        assert code == 0
        out = _json.loads(buf.getvalue())

        ids = [d.get("existing_id") for d in out["duplicates"]]
        assert "leak-id" not in ids, out
        assert out["duplicates_egress"]["withheld"] == 1

    def test_sync_json_filters_above_tier_duplicate(self, tmp_path, monkeypatch):
        import json as _json
        from contextlib import redirect_stdout

        from brain import cli

        vault = tmp_path / "vault"
        (vault / "brain").mkdir(parents=True)
        (vault / "brain" / "index.md").write_text(
            "---\nid: index\ntitle: \"Index\"\ntype: index\nclassification: Internal\n"
            "created: 2026-07-05\nupdated: 2026-07-05\n---\n\nMap.\n", encoding="utf-8",
        )
        (vault / "raw").mkdir(parents=True)

        fake_result = {
            "mode": "incremental", "added": 0, "updated": 0, "deleted": 0, "unchanged": 0,
            "chunks": 0,
            "drain": {"promoted": 0, "skipped": 0},
            "ingest": {
                "processed": [], "quarantined": [], "skipped": [], "dry_run": False,
                "duplicates": [
                    {"file": "dup.pdf", "existing_id": "leak-id", "classification": "Restricted"},
                ],
            },
        }
        monkeypatch.setattr(BrainCore, "sync", lambda self, **kw: fake_result)
        monkeypatch.setenv("BRAIN_VAULT", str(vault))
        monkeypatch.setenv("BRAIN_INDEX_DIR", str(tmp_path / "idx"))

        buf = io.StringIO()
        with redirect_stdout(buf):
            code = cli.main(["sync", "--json"])
        assert code == 0
        out = _json.loads(buf.getvalue())

        ids = [d.get("existing_id") for d in out["ingest"]["duplicates"]]
        assert "leak-id" not in ids, out
        assert out["ingest"]["duplicates_egress"]["withheld"] == 1


# ---------------------------------------------------------------------------
# S05 extra fix pass — E5: the per-file failure counter must be keyed on
# something stable across a `_claim` collision-suffix rename (e.g.
# `poison.1.pdf`), not the filename, or the count silently resets.
# ---------------------------------------------------------------------------

class TestFailureCounterKeyStability:
    def test_collision_suffixed_retry_still_accumulates_toward_quarantine(
        self, tmp_path, audit_key_env, monkeypatch
    ):
        from brain.ingest import pipeline as P
        from brain.ingest.handlers import text as text_mod

        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "poison.txt").write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        def _always_boom(cls, path):
            raise RuntimeError("simulated poison file")

        monkeypatch.setattr(text_mod.TextHandler, "extract", classmethod(_always_boom))

        processing_dir = vault / "inbox" / P.PROCESSING_DIRNAME
        processing_dir.mkdir(parents=True, exist_ok=True)
        # Force `_claim` to disambiguate: pre-occupy the natural claim slot so
        # THIS attempt's claimed copy lands as `poison.1.txt`, not `poison.txt`
        # — the exact rename that broke a name-keyed counter (E5).
        (processing_dir / "poison.txt").write_bytes(b"unrelated leftover, not the real file")

        report = core.ingest_dropzone()
        assert any(
            s["file"] == "poison.txt" and f"attempt 1/{P.MAX_INGEST_FAILURES}" in s["reason"]
            for s in report["skipped"]
        ), report

        (processing_dir / "poison.txt").unlink()  # remove the unrelated leftover

        for _ in range(P.MAX_INGEST_FAILURES - 1):
            report = core.ingest_dropzone()

        assert report["quarantined"], report
        assert report["quarantined"][0]["reason"] == "repeated_ingest_failure", (
            "a collision-suffixed retry rename must not reset the per-file "
            "failure count back to zero"
        )


# ---------------------------------------------------------------------------
# S05 extra fix pass — E6: the persisted failure counter must not outlive a
# file that eventually succeeds — a stale entry could prematurely quarantine
# a later, unrelated drop that happens to key the same.
# ---------------------------------------------------------------------------

class TestFailureCounterGarbageCollection:
    def test_success_after_transient_failures_clears_the_counter_entry(
        self, tmp_path, audit_key_env, monkeypatch
    ):
        from brain.ingest import pipeline as P
        from brain.ingest.handlers import text as text_mod

        vault = _mini_vault(tmp_path)
        (vault / "inbox" / "flaky.txt").write_text(LOREM, encoding="utf-8")
        core = _host_core(tmp_path, vault)

        real_extract = text_mod.TextHandler.extract.__func__
        calls = {"n": 0}

        def _flaky(cls, path):
            calls["n"] += 1
            if calls["n"] <= 2:
                raise RuntimeError("transient failure")
            return real_extract(cls, path)

        monkeypatch.setattr(text_mod.TextHandler, "extract", classmethod(_flaky))

        core.ingest_dropzone()
        core.ingest_dropzone()
        assert P._load_failures(vault), "counter must be non-empty after 2 failures"

        report = core.ingest_dropzone()
        assert len(report["processed"]) == 1, report
        assert P._load_failures(vault) == {}, (
            "a successful ingest must clear its failure-counter entry (E6)"
        )
